"""
Create PowerPoint Presentation for AirFly Insights Project
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_airfly_presentation():
    """Create comprehensive PowerPoint presentation for AirFly project"""
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "AirFly Insights"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(60)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 51, 102)
    title_para.alignment = PP_ALIGN.CENTER
    
    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Comprehensive Flight Delay Analysis"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(32)
    subtitle_para.font.color.rgb = RGBColor(102, 102, 102)
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # Add tagline
    tagline_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(0.6))
    tagline_frame = tagline_box.text_frame
    tagline_frame.text = "Data-Driven Insights from 5.8 Million Flight Records"
    tagline_para = tagline_frame.paragraphs[0]
    tagline_para.font.size = Pt(20)
    tagline_para.font.italic = True
    tagline_para.font.color.rgb = RGBColor(51, 102, 153)
    tagline_para.alignment = PP_ALIGN.CENTER
    
    # Slide 2: Project Overview
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Project Overview"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Comprehensive Analysis of U.S. Aviation Data"
    
    points = [
        "5,819,080 flight records analyzed",
        "14 major U.S. airline carriers",
        "300+ airports across the United States",
        "Full year of operational data (564.96 MB dataset)",
        "30+ features including delays, cancellations, routes",
        "Interactive dashboards and 40+ visualizations"
    ]
    
    for point in points:
        p = tf.add_paragraph()
        p.text = point
        p.level = 1
        p.font.size = Pt(18)
    
    # Slide 3: Key Features
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Key Features & Capabilities"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Advanced Analytics & Visualization"
    
    features = [
        ("Exploratory Data Analysis", "Comprehensive dataset exploration and quality assessment"),
        ("Feature Engineering", "Custom features: delay categories, seasons, routes, time-based metrics"),
        ("Interactive Dashboard", "Streamlit-powered web application with real-time filtering"),
        ("Geographic Visualizations", "Interactive maps showing traffic patterns and delays"),
        ("Statistical Insights", "Correlation analysis, trend detection, pattern recognition"),
        ("Performance Optimization", "Smart sampling and memory optimization for large datasets")
    ]
    
    for feature, desc in features:
        p = tf.add_paragraph()
        p.text = f"{feature}: {desc}"
        p.level = 1
        p.font.size = Pt(16)
    
    # Slide 4: Technical Stack
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Technology Stack"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    # Data Processing
    p = tf.add_paragraph()
    p.text = "Data Processing & Analysis"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    
    tools = ["Python 3.9+", "Pandas & NumPy", "Jupyter Notebook"]
    for tool in tools:
        p = tf.add_paragraph()
        p.text = tool
        p.level = 1
        p.font.size = Pt(16)
    
    # Visualization
    p = tf.add_paragraph()
    p.text = "\nVisualization & Dashboard"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    
    viz_tools = ["Matplotlib & Seaborn", "Plotly (Interactive)", "Streamlit Framework"]
    for tool in viz_tools:
        p = tf.add_paragraph()
        p.text = tool
        p.level = 1
        p.font.size = Pt(16)
    
    # Slide 5: Data Pipeline
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Data Processing Pipeline"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    pipeline_steps = [
        ("1. Data Loading", "Multiple datasets: flights, airlines, airports, weather, holidays"),
        ("2. Data Quality Assessment", "Missing value analysis, duplicate detection, outlier identification"),
        ("3. Data Preprocessing", "Missing value handling, data type optimization, filtering"),
        ("4. Feature Engineering", "Date/time features, route creation, delay categorization, seasons"),
        ("5. Memory Optimization", "Dtype optimization reducing dataset size by ~40%"),
        ("6. Analysis & Visualization", "Statistical analysis, correlation studies, comprehensive visualizations")
    ]
    
    for step, desc in pipeline_steps:
        p = tf.add_paragraph()
        p.text = step
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 51, 102)
        
        p = tf.add_paragraph()
        p.text = desc
        p.level = 1
        p.font.size = Pt(14)
    
    # Slide 6: Feature Engineering
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Feature Engineering"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Custom Features Created for Analysis"
    
    engineered_features = [
        "YEAR, MONTH, DAY - Date components for temporal analysis",
        "DAY_OF_WEEK, DAY_NAME - Weekday patterns (0-6, Monday-Sunday)",
        "DEP_HOUR - Departure hour (0-23) for hourly patterns",
        "ROUTE - Origin-Destination pair for route analysis",
        "DELAY_CATEGORY - 5-level classification (On Time to Severe)",
        "SEASON - Seasonal categorization (Winter, Spring, Summer, Fall)",
        "TOTAL_DELAY - Aggregate delay metric for comprehensive analysis"
    ]
    
    for feature in engineered_features:
        p = tf.add_paragraph()
        p.text = feature
        p.level = 1
        p.font.size = Pt(16)
    
    # Slide 7: Key Insights - Flight Patterns
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Key Insights: Flight Patterns"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    insights = [
        ("Peak Travel Period", "Summer months (June-August) show highest volume"),
        ("Busiest Day", "Friday experiences the highest flight frequency"),
        ("Peak Hours", "Morning rush: 6:00 AM - 8:00 AM"),
        ("Route Distribution", "Top routes concentrated in major hub cities"),
        ("Seasonal Variations", "Significant differences across seasons")
    ]
    
    for insight_title, insight_desc in insights:
        p = tf.add_paragraph()
        p.text = f"✈ {insight_title}"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 102, 204)
        
        p = tf.add_paragraph()
        p.text = insight_desc
        p.level = 1
        p.font.size = Pt(15)
    
    # Slide 8: Key Insights - Delay Analysis
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Key Insights: Delay Analysis"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    delay_insights = [
        ("Delay Distribution", "Most flights experience minor or no delays"),
        ("Time-Based Patterns", "Early morning flights show better performance"),
        ("Seasonal Impact", "Summer delays tend to be higher due to weather"),
        ("Distance Correlation", "Longer flights show different delay patterns"),
        ("Airline Performance", "Significant variance across carriers"),
        ("Airport Efficiency", "Hub airports show distinct delay characteristics")
    ]
    
    for insight_title, insight_desc in delay_insights:
        p = tf.add_paragraph()
        p.text = f"📊 {insight_title}"
        p.font.size = Pt(17)
        p.font.bold = True
        p.font.color.rgb = RGBColor(204, 51, 0)
        
        p = tf.add_paragraph()
        p.text = insight_desc
        p.level = 1
        p.font.size = Pt(14)
    
    # Slide 9: Visualizations
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Comprehensive Visualizations"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "40+ Charts & Interactive Maps"
    
    viz_categories = [
        ("Temporal Analysis", "Monthly trends, daily patterns, hourly distributions, seasonal variations"),
        ("Airport Analysis", "Traffic heatmaps, delay patterns, busiest airports, geographic maps"),
        ("Airline Performance", "Delays by carrier, cancellations, diversions, reliability metrics"),
        ("Route Analysis", "Top routes, distance analysis, route-specific delays"),
        ("Statistical Charts", "Distribution plots, correlation heatmaps, box plots, histograms"),
        ("Interactive Maps", "Airport delay map, route flow visualization, traffic density")
    ]
    
    for category, description in viz_categories:
        p = tf.add_paragraph()
        p.text = f"{category}: {description}"
        p.level = 1
        p.font.size = Pt(15)
    
    # Slide 10: Interactive Dashboard
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Interactive Streamlit Dashboard"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Web-Based Analytics Platform"
    
    dashboard_features = [
        "Real-time filtering by airline, airport, and date range",
        "Dynamic visualizations updating based on selections",
        "Multi-page interface: Overview, Airlines, Routes, Temporal",
        "KPI cards showing key performance indicators",
        "Export capabilities for charts and data",
        "Responsive design for various screen sizes"
    ]
    
    for feature in dashboard_features:
        p = tf.add_paragraph()
        p.text = feature
        p.level = 1
        p.font.size = Pt(18)
    
    # Slide 11: Project Structure
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Project Organization"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    structure = [
        ("📓 analysis/", "Jupyter notebooks and dashboard code"),
        ("📊 dataset/", "Flight data, airlines, airports (5.8M records)"),
        ("📚 documentation/", "Comprehensive guides and references"),
        ("🗺️ maps/", "Interactive HTML visualizations"),
        ("⚙️ configuration/", "Requirements and settings"),
        ("🧪 testing/", "Test suites and quality assurance")
    ]
    
    for folder, desc in structure:
        p = tf.add_paragraph()
        p.text = f"{folder}: {desc}"
        p.font.size = Pt(18)
        p.space_after = Pt(10)
    
    # Slide 12: Performance Optimizations
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Performance Optimizations"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Handling Large-Scale Data Efficiently"
    
    optimizations = [
        "Memory reduction by ~40% through dtype optimization",
        "Smart sampling (500K from 5.8M) for intensive computations",
        "Vectorized operations for faster processing",
        "Efficient groupby operations with optimized aggregations",
        "Incremental processing for memory-intensive tasks",
        "Caching strategies for repeated computations"
    ]
    
    for opt in optimizations:
        p = tf.add_paragraph()
        p.text = opt
        p.level = 1
        p.font.size = Pt(17)
    
    # Slide 13: Use Cases
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Applications & Use Cases"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    use_cases = [
        ("🎓 Education", "Learning data science, visualization, and feature engineering"),
        ("✈️ Aviation Industry", "Understanding operational patterns and delay causes"),
        ("💼 Business Intelligence", "Creating insights from large operational datasets"),
        ("📊 Dashboard Development", "Building interactive analytical applications"),
        ("📈 Statistical Analysis", "Applying advanced analytics to real-world data"),
        ("🤖 ML Preparation", "Feature-engineered dataset ready for predictive modeling")
    ]
    
    for use_case, description in use_cases:
        p = tf.add_paragraph()
        p.text = f"{use_case}: {description}"
        p.font.size = Pt(17)
        p.space_after = Pt(8)
    
    # Slide 14: Future Enhancements
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Future Enhancements"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Roadmap for Expansion"
    
    enhancements = [
        "Machine Learning Models - Predict flight delays using engineered features",
        "Real-time Integration - Live flight tracking and delay prediction",
        "Advanced Analytics - Clustering, anomaly detection, pattern mining",
        "Mobile Dashboard - Responsive design for mobile devices",
        "API Development - RESTful API for programmatic data access",
        "Automated Reporting - Scheduled report generation and distribution"
    ]
    
    for enhancement in enhancements:
        p = tf.add_paragraph()
        p.text = enhancement
        p.level = 1
        p.font.size = Pt(17)
    
    # Slide 15: Project Deliverables
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Project Deliverables"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    deliverables = [
        ("✅ Comprehensive Jupyter Notebook", "57 cells with complete analysis workflow"),
        ("✅ Interactive Dashboard", "Streamlit web application with multiple pages"),
        ("✅ Geographic Visualizations", "3 interactive HTML maps"),
        ("✅ Processed Dataset", "Feature-engineered data ready for analysis"),
        ("✅ Complete Documentation", "7 detailed documentation files"),
        ("✅ Test Suite", "Comprehensive testing framework"),
        ("✅ Project Presentation", "PowerPoint slides and materials")
    ]
    
    for deliverable, desc in deliverables:
        p = tf.add_paragraph()
        p.text = deliverable
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 153, 0)
        
        p = tf.add_paragraph()
        p.text = desc
        p.level = 1
        p.font.size = Pt(14)
    
    # Slide 16: Conclusion
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Conclusion"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 51, 102)
    title_para.alignment = PP_ALIGN.CENTER
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(1.5), Inches(3), Inches(7), Inches(3))
    content_frame = content_box.text_frame
    
    p = content_frame.paragraphs[0]
    p.text = "AirFly Insights demonstrates:"
    p.font.size = Pt(22)
    p.font.bold = True
    
    conclusions = [
        "Comprehensive analysis of 5.8 million flight records",
        "Advanced feature engineering and data processing",
        "Interactive visualizations and dashboard development",
        "Actionable insights for aviation industry stakeholders",
        "Scalable architecture for large-scale data analytics"
    ]
    
    for conclusion in conclusions:
        p = content_frame.add_paragraph()
        p.text = f"• {conclusion}"
        p.font.size = Pt(18)
        p.space_before = Pt(10)
    
    # Slide 17: Thank You
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Thank you text
    thanks_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
    thanks_frame = thanks_box.text_frame
    thanks_frame.text = "Thank You!"
    thanks_para = thanks_frame.paragraphs[0]
    thanks_para.font.size = Pt(60)
    thanks_para.font.bold = True
    thanks_para.font.color.rgb = RGBColor(0, 51, 102)
    thanks_para.alignment = PP_ALIGN.CENTER
    
    # Contact/info
    info_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(1.5))
    info_frame = info_box.text_frame
    
    p = info_frame.paragraphs[0]
    p.text = "AirFly Insights"
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(51, 102, 153)
    p.alignment = PP_ALIGN.CENTER
    
    p = info_frame.add_paragraph()
    p.text = "Comprehensive Flight Delay Analysis"
    p.font.size = Pt(20)
    p.font.italic = True
    p.font.color.rgb = RGBColor(102, 102, 102)
    p.alignment = PP_ALIGN.CENTER
    
    p = info_frame.add_paragraph()
    p.text = "December 2025"
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(153, 153, 153)
    p.alignment = PP_ALIGN.CENTER
    
    # Save presentation
    output_path = r"c:\Users\Kaustab das\Desktop\airfly\AirFly_Insights_Presentation.pptx"
    prs.save(output_path)
    print(f"✅ Presentation created successfully: {output_path}")
    return output_path

if __name__ == "__main__":
    create_airfly_presentation()
